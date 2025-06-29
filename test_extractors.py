#!/usr/bin/env python3
"""
Test script to verify PPTX and XLSX text extraction functionality
"""

from rag_file_monitor.text_extractors import TextExtractor
import tempfile
import os


def test_pptx_extraction():
    """Test PPTX text extraction"""
    print("Testing PPTX extraction...")

    try:
        from pptx import Presentation
        from pptx.util import Inches

        # Create a temporary PPTX file for testing
        prs = Presentation()

        # Add a slide with title and content
        slide1 = prs.slides.add_slide(prs.slide_layouts[1])
        slide1.shapes.title.text = "Test Presentation"
        slide1.shapes.placeholders[1].text = "This is a test slide with content."

        # Add another slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Second Slide"
        slide2.shapes.placeholders[1].text = "This is the second slide content."

        # Save to temporary file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs.save(tmp.name)
            tmp_path = tmp.name

        # Test extraction
        extractor = TextExtractor()
        extracted_text = extractor.extract_text(tmp_path)

        print(f"Extracted PPTX text:\n{extracted_text}")
        print(f"✓ PPTX extraction successful")

        # Cleanup
        os.unlink(tmp_path)

    except Exception as e:
        print(f"✗ PPTX extraction failed: {e}")


def test_xlsx_extraction():
    """Test XLSX text extraction"""
    print("\nTesting XLSX extraction...")

    try:
        from openpyxl import Workbook

        # Create a temporary XLSX file for testing
        wb = Workbook()
        ws = wb.active
        ws.title = "Test Sheet"

        # Add some data
        ws["A1"] = "Name"
        ws["B1"] = "Age"
        ws["C1"] = "City"
        ws["A2"] = "John Doe"
        ws["B2"] = 30
        ws["C2"] = "New York"
        ws["A3"] = "Jane Smith"
        ws["B3"] = 25
        ws["C3"] = "Los Angeles"

        # Add another sheet
        ws2 = wb.create_sheet("Second Sheet")
        ws2["A1"] = "Product"
        ws2["B1"] = "Price"
        ws2["A2"] = "Laptop"
        ws2["B2"] = 999.99

        # Save to temporary file
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            wb.save(tmp.name)
            tmp_path = tmp.name

        # Test extraction
        extractor = TextExtractor()
        extracted_text = extractor.extract_text(tmp_path)

        print(f"Extracted XLSX text:\n{extracted_text}")
        print(f"✓ XLSX extraction successful")

        # Cleanup
        os.unlink(tmp_path)

    except Exception as e:
        print(f"✗ XLSX extraction failed: {e}")


if __name__ == "__main__":
    test_pptx_extraction()
    test_xlsx_extraction()
