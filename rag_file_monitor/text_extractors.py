"""
Text extraction utilities for various file formats
"""

import logging
import chardet
from pathlib import Path

# PDF extraction
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

# DOCX extraction
try:
    from docx import Document
except ImportError:
    Document = None

# PPTX extraction
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# XLSX extraction
try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

# HTML extraction
try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None


class TextExtractor:
    """Extract text from various file formats"""

    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.fallback_text = "*** document text cannot be extracted ***"  # Default fallback text for extraction failures

    def extract_text(self, file_path: str) -> str:
        """Extract text from file based on its type"""
        path = Path(file_path)
        extension = path.suffix.lower()
        # Check if file exists and has content
        if not path.exists():
            self.logger.warning(f"File does not exist: {file_path}")
            return ""
        
        if path.stat().st_size < 10:
            self.logger.warning(f"File is too small: {file_path}")
            return ""
        try:
            result = ""
            if extension == ".pdf":
                result = self.extract_pdf(file_path)
            elif extension == ".docx":
                result = self.extract_docx(file_path)
            elif extension == ".pptx":
                result = self.extract_pptx(file_path)
            elif extension == ".xlsx":
                result = self.extract_xlsx(file_path)
            elif extension in [".html", ".htm"]:
                result = self.extract_html(file_path)
            elif extension in [".txt", ".md", ".rtf"]:
                result = self.extract_text_file(file_path)
            else:
                # Try as text file
                result = self.extract_text_file(file_path)
            if not result.strip():
                self.logger.warning(f"No text extracted from {file_path}, using fallback text")
                return self.fallback_text
            return result

        except Exception as e:
            self.logger.warning(f"Error extracting text from {file_path}: {str(e)}")
            return self.fallback_text

    def extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        if PyPDF2 is None:
            raise ImportError("PyPDF2 is required for PDF extraction")

        text = ""
        try:
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            self.logger.error(f"Error reading PDF {file_path}: {str(e)}")

        return text.strip()

    def extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        if Document is None:
            raise ImportError("python-docx is required for DOCX extraction")

        try:
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return "\n".join(text)
        except Exception as e:
            self.logger.error(f"Error reading DOCX {file_path}: {str(e)}")
            return ""

    def extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file"""
        if Presentation is None:
            raise ImportError("python-pptx is required for PPTX extraction")

        try:
            prs = Presentation(file_path)
            text = []

            for slide in prs.slides:
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():  # Only add non-empty text
                            text.append(shape.text.strip())

                    # Handle tables if present
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text.append(" | ".join(row_text))

                # Add separator between slides
                if text and not text[-1].startswith("--- Slide"):
                    text.append(f"--- Slide {len([t for t in text if t.startswith('--- Slide')]) + 1} ---")

            return "\n".join(text)

        except Exception as e:
            self.logger.error(f"Error reading PPTX {file_path}: {str(e)}")
            return ""

    def extract_xlsx(self, file_path: str) -> str:
        """Extract text from XLSX file"""
        if load_workbook is None:
            raise ImportError("openpyxl is required for XLSX extraction")

        try:
            workbook = load_workbook(file_path, read_only=True, data_only=True)
            text = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"--- Sheet: {sheet_name} ---")

                # Get the used range to avoid empty cells
                if sheet.max_row > 0 and sheet.max_column > 0:
                    for row in sheet.iter_rows(
                        min_row=1, max_row=sheet.max_row, min_col=1, max_col=sheet.max_column, values_only=True
                    ):
                        # Filter out None values and convert to strings
                        row_values = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                        if row_values:  # Only add non-empty rows
                            text.append(" | ".join(row_values))

                text.append("")  # Add blank line between sheets

            workbook.close()
            return "\n".join(text)

        except Exception as e:
            self.logger.error(f"Error reading XLSX {file_path}: {str(e)}")
            return ""

    def extract_html(self, file_path: str) -> str:
        """Extract text from HTML file"""
        if BeautifulSoup is None:
            raise ImportError("beautifulsoup4 is required for HTML extraction")

        try:
            with open(file_path, "r", encoding="utf-8") as file:
                content = file.read()

            soup = BeautifulSoup(content, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Get text and clean it up
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = "\n".join(chunk for chunk in chunks if chunk)

            return text

        except Exception as e:
            self.logger.error(f"Error reading HTML {file_path}: {str(e)}")
            return ""

    def extract_text_file(self, file_path: str) -> str:
        """Extract text from plain text file with memory optimization"""
        try:
            # Detect encoding first with limited data
            with open(file_path, "rb") as file:
                # Read only a sample for encoding detection to save memory
                raw_data = file.read(10000)  # Read first 10KB for encoding detection

            result = chardet.detect(raw_data)
            encoding = result["encoding"] or "utf-8"

            # For very large files, implement size limit
            with open(file_path, "r", encoding=encoding, errors="ignore") as file:
                # Read in chunks for very large files
                max_size = 10 * 1024 * 1024  # 10MB limit
                content = file.read(max_size)

                # Check if file was truncated
                if len(content) == max_size:
                    self.logger.warning(f"Large file {file_path} was truncated to {max_size} bytes")

            return content

        except Exception as e:
            self.logger.error(f"Error reading text file {file_path}: {str(e)}")
            return ""
