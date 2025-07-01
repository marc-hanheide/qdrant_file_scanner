"""
Unit tests for text extractors module
"""

import pytest
import tempfile
import os
from pathlib import Path
import sys

# Add the project root to sys.path
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

from rag_file_monitor.text_extractors import TextExtractor


def _has_docx_support():
    """Check if python-docx is available"""
    try:
        import docx
        return True
    except ImportError:
        return False


def _has_pptx_support():
    """Check if python-pptx is available"""
    try:
        import pptx
        return True
    except ImportError:
        return False


def _has_xlsx_support():
    """Check if openpyxl is available"""
    try:
        import openpyxl
        return True
    except ImportError:
        return False


class TestTextExtractor:
    """Test cases for TextExtractor class"""

    def setup_method(self):
        """Set up test fixtures"""
        self.extractor = TextExtractor()

    def test_extract_plain_text(self):
        """Test extraction of plain text files"""
        test_content = "This is a test text file.\nWith multiple lines.\nAnd some content."
        
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as tmp:
            tmp.write(test_content)
            tmp_path = tmp.name

        try:
            extracted = self.extractor.extract_text(tmp_path)
            assert extracted == test_content
        finally:
            os.unlink(tmp_path)

    def test_extract_markdown(self):
        """Test extraction of markdown files"""
        test_content = "# Test Header\n\nThis is markdown content.\n\n- List item 1\n- List item 2"
        
        with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False) as tmp:
            tmp.write(test_content)
            tmp_path = tmp.name

        try:
            extracted = self.extractor.extract_text(tmp_path)
            assert extracted == test_content
        finally:
            os.unlink(tmp_path)

    def test_extract_html(self):
        """Test extraction of HTML files"""
        html_content = """<html>
        <head><title>Test Page</title></head>
        <body>
            <h1>Test Header</h1>
            <p>This is HTML content.</p>
            <ul>
                <li>List item 1</li>
                <li>List item 2</li>
            </ul>
        </body>
        </html>"""
        
        expected_text_parts = ["Test Page", "Test Header", "This is HTML content.", "List item 1", "List item 2"]
        
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False) as tmp:
            tmp.write(html_content)
            tmp_path = tmp.name

        try:
            extracted = self.extractor.extract_text(tmp_path)
            for part in expected_text_parts:
                assert part in extracted
        finally:
            os.unlink(tmp_path)

    @pytest.mark.skipif(not _has_docx_support(), reason="python-docx not available")
    def test_extract_docx(self):
        """Test extraction of DOCX files"""
        try:
            from docx import Document
            
            # Create a temporary DOCX file
            doc = Document()
            doc.add_heading('Test Document', 0)
            doc.add_paragraph('This is a test paragraph.')
            doc.add_paragraph('This is another paragraph.')
            
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                doc.save(tmp.name)
                tmp_path = tmp.name

            try:
                extracted = self.extractor.extract_text(tmp_path)
                assert "Test Document" in extracted
                assert "This is a test paragraph." in extracted
                assert "This is another paragraph." in extracted
            finally:
                os.unlink(tmp_path)
        except ImportError:
            pytest.skip("python-docx not available")

    @pytest.mark.skipif(not _has_pptx_support(), reason="python-pptx not available")
    def test_extract_pptx(self):
        """Test extraction of PPTX files"""
        try:
            from pptx import Presentation
            
            # Create a temporary PPTX file
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Test Presentation"
            slide.shapes.placeholders[1].text = "This is test content."
            
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                prs.save(tmp.name)
                tmp_path = tmp.name

            try:
                extracted = self.extractor.extract_text(tmp_path)
                assert "Test Presentation" in extracted
                assert "This is test content." in extracted
            finally:
                os.unlink(tmp_path)
        except ImportError:
            pytest.skip("python-pptx not available")

    @pytest.mark.skipif(not _has_xlsx_support(), reason="openpyxl not available")
    def test_extract_xlsx(self):
        """Test extraction of XLSX files"""
        try:
            from openpyxl import Workbook
            
            # Create a temporary XLSX file
            wb = Workbook()
            ws = wb.active
            ws.title = "Test Sheet"
            ws["A1"] = "Name"
            ws["B1"] = "Age"
            ws["A2"] = "John Doe"
            ws["B2"] = 30
            
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
                wb.save(tmp.name)
                tmp_path = tmp.name

            try:
                extracted = self.extractor.extract_text(tmp_path)
                assert "Name" in extracted
                assert "Age" in extracted
                assert "John Doe" in extracted
                assert "30" in extracted
            finally:
                os.unlink(tmp_path)
        except ImportError:
            pytest.skip("openpyxl not available")

    def test_extract_nonexistent_file(self):
        """Test extraction of non-existent file"""
        extracted = self.extractor.extract_text("/nonexistent/file.txt")
        assert extracted == ""

    def test_extract_unsupported_format(self):
        """Test extraction of unsupported file format"""
        with tempfile.NamedTemporaryFile(suffix='.unknown', delete=False) as tmp:
            tmp.write(b"Some binary content")
            tmp_path = tmp.name

        try:
            # Should fall back to text extraction
            extracted = self.extractor.extract_text(tmp_path)
            # Might be empty if binary content can't be decoded
            assert isinstance(extracted, str)
        finally:
            os.unlink(tmp_path)


if __name__ == "__main__":
    pytest.main([__file__])
